"""
gen_slides.py — Tạo slide thuyết trình Digital Clock SN32F407
Dùng python-pptx, thiết kế sạch: nền trắng / navy, font Calibri
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import copy, os

OUT = r"C:\Sonix\Clock_SN32F407\Slide_Digital_Clock_SN32F407.pptx"

# ─── Màu chủ đạo ────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x0D, 0x2F, 0x5E)   # nền tiêu đề
BLUE    = RGBColor(0x15, 0x5E, 0xA8)   # accent
CYAN    = RGBColor(0x00, 0xB0, 0xD8)   # highlight
GREEN   = RGBColor(0x1A, 0x8A, 0x3C)   # success / good
ORANGE  = RGBColor(0xD4, 0x6B, 0x08)   # warning / note
RED     = RGBColor(0xC0, 0x20, 0x20)   # error / bad
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY   = RGBColor(0xF2, 0xF4, 0xF8)   # nền content nhạt
DGRAY   = RGBColor(0x44, 0x44, 0x55)   # text phụ
BLACK   = RGBColor(0x11, 0x11, 0x22)

prs = Presentation()
prs.slide_width  = Inches(13.33)   # 16:9
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height
BLANK = prs.slide_layouts[6]       # hoàn toàn trắng

# ════════════════════════════════════════════════════════════════════════════
# Helpers
# ════════════════════════════════════════════════════════════════════════════
def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE=1
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        if line_w: shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             size=20, bold=False, italic=False, color=BLACK,
             align=PP_ALIGN.LEFT, wrap=True, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = font
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return tb

def add_para(tf, text, size=16, bold=False, color=BLACK,
             align=PP_ALIGN.LEFT, font="Calibri", italic=False, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text  = text
    run.font.name   = font
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return p

def add_textbox_multi(slide, lines, x, y, w, h, default_size=16, wrap=True):
    """lines = list of (text, size, bold, color, align, italic)"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    first = True
    for item in lines:
        text  = item[0]
        size  = item[1] if len(item)>1 else default_size
        bold  = item[2] if len(item)>2 else False
        color = item[3] if len(item)>3 else BLACK
        align = item[4] if len(item)>4 else PP_ALIGN.LEFT
        italic= item[5] if len(item)>5 else False
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name  = "Calibri"
        run.font.size  = Pt(size)
        run.font.bold  = bold
        run.font.color.rgb = color
        run.font.italic = italic
    return tb

def header_bar(slide, title, subtitle=None):
    """Thanh tiêu đề navy ở trên"""
    bar_h = Inches(1.2)
    add_rect(slide, 0, 0, W, bar_h, fill=NAVY)
    # Đường accent cyan dưới bar
    add_rect(slide, 0, bar_h, W, Pt(4).emu, fill=CYAN)
    add_text(slide, title, Inches(0.35), Inches(0.12), Inches(10), Inches(0.7),
             size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, Inches(0.38), Inches(0.78), Inches(10), Inches(0.38),
                 size=16, italic=True, color=CYAN, align=PP_ALIGN.LEFT)

def slide_number(slide, n, total):
    add_text(slide, f"{n} / {total}", W - Inches(1.1), H - Inches(0.38),
             Inches(0.9), Inches(0.3), size=11, color=DGRAY, align=PP_ALIGN.RIGHT)

def bullet_box(slide, items, x, y, w, h, title=None, title_color=BLUE,
               bg=LGRAY, bullet="▸", size=15, title_size=17):
    add_rect(slide, x, y, w, h, fill=bg)
    tb = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.12),
                                   w - Inches(0.36), h - Inches(0.24))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0]; first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run(); run.text = title
        run.font.name = "Calibri"; run.font.size = Pt(title_size)
        run.font.bold = True; run.font.color.rgb = title_color

    for item in items:
        if first: p = tf.paragraphs[0]; first = False
        else: p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        if isinstance(item, tuple):
            run.text = bullet + " " + item[0]
            run.font.color.rgb = item[1] if len(item)>1 else BLACK
            run.font.bold = item[2] if len(item)>2 else False
        else:
            run.text = bullet + " " + item
            run.font.color.rgb = BLACK
        run.font.name = "Calibri"; run.font.size = Pt(size)

def code_box(slide, code_text, x, y, w, h, title=None):
    add_rect(slide, x, y, w, h, fill=RGBColor(0x1E,0x1E,0x2E))
    if title:
        add_text(slide, title, x+Inches(0.1), y+Inches(0.05), w, Inches(0.3),
                 size=12, bold=True, color=CYAN)
        yo = Inches(0.3)
    else:
        yo = Inches(0.1)
    tb = slide.shapes.add_textbox(x+Inches(0.15), y+yo, w-Inches(0.3), h-yo-Inches(0.1))
    tf = tb.text_frame; tf.word_wrap = False
    first = True
    for line in code_text.split('\n'):
        if first: p = tf.paragraphs[0]; first = False
        else: p = tf.add_paragraph()
        run = p.add_run(); run.text = line
        run.font.name = "Courier New"; run.font.size = Pt(11.5)
        if line.strip().startswith('//') or line.strip().startswith('/*') or line.strip().startswith('*'):
            run.font.color.rgb = RGBColor(0x6A,0xC0,0x6A)
        elif any(kw in line for kw in ['void ','uint','int ','if ','while','return','#define','volatile']):
            run.font.color.rgb = RGBColor(0x56,0x9C,0xD6)
        elif '"' in line or "'" in line:
            run.font.color.rgb = RGBColor(0xCE,0x91,0x78)
        else:
            run.font.color.rgb = RGBColor(0xD4,0xD4,0xD4)

TOTAL = 18

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
# Full navy background
add_rect(s, 0, 0, W, H, fill=NAVY)
# Stripe cyan trên
add_rect(s, 0, 0, W, Inches(0.18), fill=CYAN)
# Stripe cyan dưới
add_rect(s, 0, H - Inches(0.18), W, Inches(0.18), fill=CYAN)
# Dải accent giữa
add_rect(s, 0, Inches(2.7), W, Pt(3).emu, fill=CYAN)
add_rect(s, 0, Inches(4.6), W, Pt(3).emu, fill=RGBColor(0x33,0x77,0xBB))

# Icon đồng hồ giả (hình tròn)
add_rect(s, Inches(0.5), Inches(1.5), Inches(1.8), Inches(1.8),
         fill=RGBColor(0x1A,0x4F,0x9A), line=CYAN, line_w=Pt(2).emu)
add_text(s, "🕐", Inches(0.62), Inches(1.6), Inches(1.5), Inches(1.5),
         size=48, color=CYAN, align=PP_ALIGN.CENTER)

# Tiêu đề chính
add_text(s, "DIGITAL CLOCK", Inches(2.5), Inches(1.1), Inches(10), Inches(1.1),
         size=52, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(s, "Đồng Hồ Kỹ Thuật Số", Inches(2.5), Inches(2.1), Inches(10), Inches(0.6),
         size=26, italic=True, color=CYAN, align=PP_ALIGN.LEFT)

# Thông tin chip
add_text(s, "SN32F407  |  ARM® Cortex®-M0  |  12 MHz IHRC",
         Inches(2.5), Inches(2.85), Inches(10), Inches(0.55),
         size=20, color=RGBColor(0xAA,0xCC,0xFF), align=PP_ALIGN.LEFT)

# Divider info
add_text(s, "Keil MDK Simulator  ·  LED 7-Segment  ·  I2C EEPROM  ·  PWM Buzzer  ·  WDT",
         Inches(0.4), Inches(4.8), Inches(12.5), Inches(0.5),
         size=16, color=RGBColor(0x88,0xAA,0xDD), align=PP_ALIGN.CENTER)

add_text(s, "TRƯỜNG ĐẠI HỌC KỸ THUẬT  |  Môn: Hệ Thống Nhúng  |  2025–2026",
         Inches(0.4), Inches(6.7), Inches(12.5), Inches(0.45),
         size=13, color=RGBColor(0x66,0x88,0xAA), align=PP_ALIGN.CENTER)

slide_number(s, 1, TOTAL)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Mục lục
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "Nội Dung Thuyết Trình", "Agenda")
slide_number(s, 2, TOTAL)

topics = [
    ("01", "Giới thiệu đề bài & yêu cầu"),
    ("02", "Kiến trúc phần cứng — SN32F407"),
    ("03", "Sơ đồ khối hệ thống"),
    ("04", "State Machine — Trái tim của chương trình"),
    ("05", "Timer & SysTick — Cơ sở thời gian"),
    ("06", "LED 7 Đoạn — Hiển thị Multiplexing"),
    ("07", "I2C & EEPROM — Lưu trữ không mất dữ liệu"),
    ("08", "Buzzer PWM — Báo thức bằng phần cứng"),
    ("09", "Watchdog Timer & An toàn hệ thống"),
    ("10", "Demo trên Keil Simulator"),
    ("11", "Kết quả & Đánh giá"),
    ("12", "Hỏi & Đáp"),
]
col1 = topics[:6]
col2 = topics[6:]
cx = Inches(0.4)
cy = Inches(1.45)
for i, (num, topic) in enumerate(col1):
    y = cy + Inches(i * 0.83)
    add_rect(s, cx, y, Inches(0.55), Inches(0.6), fill=BLUE)
    add_text(s, num, cx, y, Inches(0.55), Inches(0.6),
             size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, topic, cx+Inches(0.62), y+Inches(0.1), Inches(5.6), Inches(0.5),
             size=16, color=BLACK)

cx2 = Inches(6.8)
for i, (num, topic) in enumerate(col2):
    y = cy + Inches(i * 0.83)
    add_rect(s, cx2, y, Inches(0.55), Inches(0.6), fill=NAVY)
    add_text(s, num, cx2, y, Inches(0.55), Inches(0.6),
             size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, topic, cx2+Inches(0.62), y+Inches(0.1), Inches(5.6), Inches(0.5),
             size=16, color=BLACK)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Yêu cầu đề bài
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "Yêu Cầu Đề Bài", "15 tính năng cần thực hiện")
slide_number(s, 3, TOTAL)

reqs_left = [
    "Hiển thị giờ:phút trên LED 7 đoạn 4 chữ số",
    "Đếm giờ chính xác — dựa trên SysTick 1ms",
    "Cài đặt giờ bằng phím SETUP + PLUS/MINUS",
    "Cài đặt phút bằng phím SETUP + PLUS/MINUS",
    "Chữ số đang cài nhấp nháy 1Hz",
    "Timeout 30s tự thoát khỏi SET mode",
    "Wrap-around: 23→0, 59→0",
    "LED D0 (P3.8) nhấp nháy mỗi giây",
]
reqs_right = [
    "Cài đặt giờ báo thức (SETUP+ALARM)",
    "Cài đặt phút báo thức",
    "Lưu báo thức vào EEPROM AT24C02",
    "Báo thức kêu 5 giây khi đúng giờ",
    "Pip 300ms xác nhận khi nhấn phím",
    "Nhấn phím tắt còi báo thức",
    "Watchdog Timer — tự reset nếu firmware treo",
]

bullet_box(s, reqs_left,  Inches(0.3), Inches(1.38), Inches(6.1), Inches(5.6),
           title="✅ Tính năng đồng hồ & giao diện", title_color=GREEN, bullet="✓", size=14.5)
bullet_box(s, reqs_right, Inches(6.7), Inches(1.38), Inches(6.3), Inches(5.6),
           title="✅ Tính năng báo thức & hệ thống",  title_color=BLUE,  bullet="✓", size=14.5)

add_text(s, "→ TẤT CẢ 15 YÊU CẦU ĐÃ ĐƯỢC THỰC HIỆN",
         Inches(0.3), Inches(6.9), Inches(13), Inches(0.45),
         size=15, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Kiến trúc phần cứng
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "Kiến Trúc Phần Cứng", "SN32F407 — ARM Cortex-M0 @ 12 MHz")
slide_number(s, 4, TOTAL)

# Bảng thông số MCU
specs = [
    ("Lõi CPU",          "ARM Cortex-M0, 32-bit RISC"),
    ("Clock",            "IHRC 12 MHz (nội bộ, ±5%)"),
    ("Flash / RAM",      "64 KB / 8 KB"),
    ("Timer 16-bit",     "CT16B0, CT16B1 — PWM & đếm"),
    ("I2C",              "I2C0 @ P0.10/P0.11 — giao tiếp EEPROM"),
    ("GPIO",             "P0, P1, P3 — segment, COM, LED, phím"),
    ("Watchdog",         "WDT nội — timeout ~250ms"),
    ("Package",          "LQFP-64"),
]

add_rect(s, Inches(0.3), Inches(1.38), Inches(6.2), Inches(5.7), fill=LGRAY)
add_text(s, "Thông số SN32F407", Inches(0.45), Inches(1.45), Inches(6), Inches(0.42),
         size=16, bold=True, color=NAVY)

for i, (k, v) in enumerate(specs):
    y = Inches(1.95 + i*0.6)
    bg = WHITE if i%2==0 else LGRAY
    add_rect(s, Inches(0.32), y, Inches(6.16), Inches(0.57), fill=bg)
    add_text(s, k, Inches(0.42), y+Inches(0.08), Inches(2.0), Inches(0.42),
             size=13.5, bold=True, color=NAVY)
    add_text(s, v, Inches(2.55), y+Inches(0.08), Inches(3.9), Inches(0.42),
             size=13.5, color=BLACK)

# Sơ đồ đơn giản bên phải
rx = Inches(6.9)
add_rect(s, rx, Inches(1.38), Inches(6.1), Inches(5.7), fill=LGRAY)
add_text(s, "Kết nối ngoại vi", rx+Inches(0.15), Inches(1.45),
         Inches(5.8), Inches(0.4), size=16, bold=True, color=NAVY)

periph = [
    ("P0.10 / P0.11",  "→",  "I2C0 → AT24C02 EEPROM"),
    ("P0.0 – P0.7",    "→",  "GPIO0 → SEG A–H (LED 7 đoạn)"),
    ("P1.9 – P1.12",   "→",  "GPIO1 → COM0–COM3 (quét)"),
    ("P3.0",           "→",  "CT16B0 PWM0 → Buzzer piezo"),
    ("P3.8",           "→",  "GPIO3 → LED D0 (nhấp nháy)"),
    ("PA.0–PA.3",      "→",  "GPIO → Phím SETUP/PLUS/MINUS/ALARM"),
    ("SysTick",        "→",  "RELOAD=11999, ngắt 1ms"),
    ("CT16B1 MR9",     "→",  "11999 → timer 1ms (MR9 match)"),
    ("CT16B0 MR9/MR0", "→",  "19999/9999 → buzzer 600Hz 50%"),
]
for i, (pin, arr, desc) in enumerate(periph):
    y = Inches(1.92 + i*0.57)
    bg = WHITE if i%2==0 else LGRAY
    add_rect(s, rx+Inches(0.02), y, Inches(6.06), Inches(0.54), fill=bg)
    add_text(s, pin,  rx+Inches(0.12), y+Inches(0.07), Inches(1.9), Inches(0.4),
             size=12.5, bold=True, color=BLUE, font="Courier New")
    add_text(s, arr,  rx+Inches(2.05), y+Inches(0.07), Inches(0.3), Inches(0.4),
             size=12.5, color=DGRAY)
    add_text(s, desc, rx+Inches(2.42), y+Inches(0.07), Inches(3.55), Inches(0.4),
             size=12.5, color=BLACK)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Sơ đồ khối hệ thống
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "Sơ Đồ Khối Hệ Thống", "Luồng dữ liệu từ phím bấm đến hiển thị")
slide_number(s, 5, TOTAL)

def block(slide, label, sub, x, y, w=Inches(1.9), h=Inches(0.85),
          fg=WHITE, bg=BLUE):
    add_rect(slide, x, y, w, h, fill=bg, line=WHITE, line_w=Pt(1).emu)
    add_text(slide, label, x, y+Inches(0.04), w, Inches(0.48),
             size=13, bold=True, color=fg, align=PP_ALIGN.CENTER)
    if sub:
        add_text(slide, sub, x, y+Inches(0.44), w, Inches(0.35),
                 size=10, italic=True, color=RGBColor(0xCC,0xEE,0xFF)
                 if fg==WHITE else DGRAY, align=PP_ALIGN.CENTER)

def arrow(slide, x1, y1, x2, y2, label=""):
    from pptx.util import Emu
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = NAVY
    line.line.width = Pt(1.8).emu
    if label:
        mx = (x1+x2)//2; my = (y1+y2)//2
        add_text(slide, label, mx-Inches(0.5), my-Inches(0.2),
                 Inches(1), Inches(0.3), size=10, color=DGRAY, align=PP_ALIGN.CENTER)

# Row 1 — Input
block(s, "Phím Bấm", "KEY_SCAN()", Inches(0.2),  Inches(1.55), bg=NAVY)
block(s, "SysTick", "1ms ISR",     Inches(5.5),  Inches(1.55), bg=NAVY)
block(s, "EEPROM", "AT24C02 I2C",  Inches(10.8), Inches(1.55), bg=NAVY)

# Arrow vào main
arrow(s, Inches(1.15), Inches(2.0), Inches(4.8), Inches(3.2))
arrow(s, Inches(6.45), Inches(2.0), Inches(6.45), Inches(3.1))
arrow(s, Inches(11.75), Inches(2.0), Inches(8.1), Inches(3.2))

# Row 2 — CPU
block(s, "main() Loop", "while(1)", Inches(4.8), Inches(3.1), w=Inches(2.5), bg=BLUE)

# Arrow ra các output
arrow(s, Inches(5.0), Inches(3.95), Inches(1.5), Inches(5.0))
arrow(s, Inches(6.05), Inches(3.95), Inches(6.05), Inches(5.0))
arrow(s, Inches(7.1), Inches(3.95), Inches(10.5), Inches(5.0))

# Row 3 — Output
block(s, "State Machine", "process_key()", Inches(0.2),  Inches(5.0), bg=RGBColor(0x0A,0x6A,0x40))
block(s, "LED 7-Seg", "Digital_Scan()",   Inches(5.1),  Inches(5.0), bg=RGBColor(0x0A,0x6A,0x40))
block(s, "Buzzer PWM", "CT16B0",          Inches(10.0), Inches(5.0), bg=RGBColor(0x0A,0x6A,0x40))

# WDT nhỏ ở góc
block(s, "WDT Feed", "mỗi vòng lặp", Inches(4.9), Inches(5.0),
      w=Inches(1.8), h=Inches(0.75), bg=ORANGE)

add_text(s, "HH:MM → LED", Inches(5.1), Inches(5.9), Inches(1.9), Inches(0.3),
         size=10, italic=True, color=DGRAY, align=PP_ALIGN.CENTER)
add_text(s, "600Hz tone", Inches(10.0), Inches(5.9), Inches(1.9), Inches(0.3),
         size=10, italic=True, color=DGRAY, align=PP_ALIGN.CENTER)

# LED D0
block(s, "LED D0", "P3.8 blink", Inches(2.4), Inches(5.0),
      w=Inches(1.5), h=Inches(0.75), bg=RGBColor(0x8B,0x00,0x8B))

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — State Machine
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "State Machine — Trái Tim Chương Trình",
           "5 trạng thái, chuyển đổi bởi phím SETUP và ALARM")
slide_number(s, 6, TOTAL)

# Vẽ 5 state dạng vòng tròn giả (rect bo góc không có sẵn → dùng rect)
states = [
    ("NORMAL\n(0)",         Inches(5.4),  Inches(2.0),  BLUE),
    ("SET_HOUR\n(1)",       Inches(1.5),  Inches(3.8),  GREEN),
    ("SET_MIN\n(2)",        Inches(3.2),  Inches(5.6),  GREEN),
    ("SET_ALARM_HOUR\n(3)", Inches(7.8),  Inches(3.8),  ORANGE),
    ("SET_ALARM_MIN\n(4)",  Inches(6.2),  Inches(5.6),  ORANGE),
]
for name, x, y, color in states:
    add_rect(s, x, y, Inches(2.2), Inches(0.82), fill=color, line=WHITE, line_w=Pt(1.5).emu)
    add_text(s, name, x, y+Inches(0.05), Inches(2.2), Inches(0.72),
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Chú thích chuyển state
transitions = [
    (Inches(5.9), Inches(1.45), "Khởi động"),
    (Inches(2.5),  Inches(3.5),  "SETUP →"),
    (Inches(3.9),  Inches(5.2),  "SETUP →"),
    (Inches(1.0),  Inches(5.2),  "← SETUP (timeout)"),
    (Inches(8.5),  Inches(3.5),  "← ALARM"),
    (Inches(6.9),  Inches(5.2),  "← ALARM"),
    (Inches(9.4),  Inches(5.2),  "ALARM (timeout) →"),
]

# Bảng chuyển state
rows = [
    ("NORMAL",         "SETUP",    "SET_HOUR"),
    ("SET_HOUR",       "SETUP",    "SET_MIN"),
    ("SET_MIN",        "SETUP",    "NORMAL"),
    ("NORMAL",         "ALARM",    "SET_ALARM_HOUR"),
    ("SET_ALARM_HOUR", "ALARM",    "SET_ALARM_MIN"),
    ("SET_ALARM_MIN",  "ALARM",    "NORMAL"),
    ("Bất kỳ SET",     "Timeout 30s", "NORMAL"),
]
add_rect(s, Inches(9.5), Inches(1.38), Inches(3.6), Inches(5.7), fill=LGRAY)
add_text(s, "Bảng chuyển trạng thái", Inches(9.6), Inches(1.44),
         Inches(3.4), Inches(0.4), size=14, bold=True, color=NAVY)

hdrs = ["Từ", "Phím", "Đến"]
col_w = [1.15, 1.05, 1.25]
col_x = [9.62, 10.77, 11.82]
for ci, h in enumerate(hdrs):
    add_rect(s, Inches(col_x[ci]), Inches(1.88),
             Inches(col_w[ci]), Inches(0.38), fill=NAVY)
    add_text(s, h, Inches(col_x[ci]), Inches(1.9),
             Inches(col_w[ci]), Inches(0.34), size=12, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)

for ri, (fr, key, to) in enumerate(rows):
    bg = WHITE if ri%2==0 else LGRAY
    y = Inches(2.28 + ri*0.58)
    for ci, val in enumerate([fr, key, to]):
        add_rect(s, Inches(col_x[ci]), y, Inches(col_w[ci]), Inches(0.55), fill=bg)
        col = GREEN if val.startswith("NORMAL") else (BLUE if val.startswith("SET_") and "ALARM" not in val else (ORANGE if "ALARM" in val else BLACK))
        add_text(s, val, Inches(col_x[ci]+0.04), y+Inches(0.07),
                 Inches(col_w[ci]-0.08), Inches(0.41),
                 size=11, color=col, bold=val in ["NORMAL","SET_HOUR","SET_MIN","SET_ALARM_HOUR","SET_ALARM_MIN"],
                 align=PP_ALIGN.CENTER)

add_text(s, "PLUS/MINUS: thay đổi giá trị trong SET state",
         Inches(0.3), Inches(7.05), Inches(9.0), Inches(0.35),
         size=13, italic=True, color=DGRAY)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Timer & SysTick
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "Cơ Sở Thời Gian — SysTick & CT16B1",
           "Tạo tick 1ms — nền tảng của mọi timing trong chương trình")
slide_number(s, 7, TOTAL)

# Công thức
add_rect(s, Inches(0.3), Inches(1.4), Inches(12.7), Inches(1.3),
         fill=RGBColor(0xE8,0xF4,0xFF), line=BLUE, line_w=Pt(1.5).emu)
add_textbox_multi(s, [
    ("Công thức:  RELOAD = (HCLK ÷ Tần_số_ngắt) − 1", 17, True, NAVY),
    ("           = (12,000,000 ÷ 1,000) − 1  =  11,999     →  Ngắt mỗi 1ms", 17, False, BLUE),
], Inches(0.5), Inches(1.48), Inches(12.3), Inches(1.1))

# Code box SysTick
code_box(s, """// SysTick_1ms.c
SysTick->CTRL = 0;               // Tắt trước khi cấu hình
SysTick->LOAD = 11999;           // RELOAD = 12MHz/1000 - 1
SysTick->VAL  = 0;               // Reset bộ đếm
SysTick->CTRL = 0x07;            // CLKSOURCE|TICKINT|ENABLE

void SysTick_Handler(void) {
    timer_1ms_flag = 1;          // Chỉ set cờ — KHÔNG xử lý logic ở đây
}""",
         Inches(0.3), Inches(2.85), Inches(6.1), Inches(3.0), title="SysTick (dùng trong Simulator)")

code_box(s, """// CT16B1.c — dùng trên phần cứng thật
SN_CT16B1->MR9   = 12*1000 - 1; // = 11999 → 1ms
SN_CT16B1->MCTRL = (1<<30)|(1<<29); // MR9IE + MR9RST

void CT16B1_IRQHandler(void) {
    uint32_t ris = SN_CT16B1->RIS;
    if (ris & (1<<5)) {          // MR9 match?
        timer_1ms_flag = 1;
    }
    SN_CT16B1->IC = ris;         // Clear interrupt flag
}""",
         Inches(6.7), Inches(2.85), Inches(6.3), Inches(3.0), title="CT16B1 (phần cứng thật)")

# Timing ladder
add_rect(s, Inches(0.3), Inches(5.95), Inches(12.7), Inches(1.32), fill=LGRAY)
add_textbox_multi(s, [
    ("Luồng timing:", 14, True, NAVY),
    ("timer_1ms_flag=1 (ISR)  →  main loop đọc flag  →  Digital_Scan() + tick_second() + update_buzzer()", 14, False, BLACK),
    ("Sau 1000ms: g_sec++ → sau 60s: g_min++ → sau 60min: g_hour++ → wrap 23→0", 13, False, BLUE),
], Inches(0.45), Inches(5.98), Inches(12.4), Inches(1.2))

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — LED 7 Đoạn Multiplexing
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "LED 7 Đoạn — Multiplexing 250Hz",
           "4 chữ số, quét tuần tự, 1ms/digit → 250Hz refresh rate")
slide_number(s, 8, TOTAL)

code_box(s, """void Digital_Scan(void) {
    // Bước 1: Tắt segment (tránh ghost)
    SN_GPIO0->BCLR = 0xFF;
    // Bước 2: Tắt COM hiện tại
    SN_GPIO1->BSET = (1 << (9 + com_scan)); // COM active LOW → set = tắt
    // Bước 3: Chuyển sang COM tiếp theo
    com_scan = (com_scan + 1) % 4;
    // Bước 4: Bật COM mới
    SN_GPIO1->BCLR = (1 << (9 + com_scan));
    // Bước 5: Xuất segment data
    SN_GPIO0->BSET = segment_buff[com_scan];
}""",
         Inches(0.3), Inches(1.42), Inches(6.3), Inches(3.3), title="Digital_Scan() — gọi mỗi 1ms")

# Bảng SEGMENT_TABLE
add_rect(s, Inches(6.8), Inches(1.42), Inches(6.2), Inches(3.3), fill=LGRAY)
add_text(s, "SEGMENT_TABLE[]", Inches(6.9), Inches(1.48), Inches(6.0), Inches(0.4),
         size=15, bold=True, color=NAVY)
seg_data = [
    ("0", "0x3F", "0b0011_1111", "a b c d e f"),
    ("1", "0x06", "0b0000_0110", "b c"),
    ("2", "0x5B", "0b0101_1011", "a b d e g"),
    ("3", "0x4F", "0b0100_1111", "a b c d g"),
    ("4", "0x66", "0b0110_0110", "b c f g"),
    ("5", "0x6D", "0b0110_1101", "a c d f g"),
    ("8", "0x7F", "0b0111_1111", "all on"),
    ("-", "0x40", "0b0100_0000", "g only (dash)"),
]
hd = ["Digit","Hex","Binary","Segments ON"]
cw = [0.55, 0.75, 1.4, 1.5]
cx2 = [6.85, 7.42, 8.19, 9.61]
for ci, h in enumerate(hd):
    add_rect(s, Inches(cx2[ci]), Inches(1.92), Inches(cw[ci]), Inches(0.36), fill=NAVY)
    add_text(s, h, Inches(cx2[ci]), Inches(1.94), Inches(cw[ci]), Inches(0.3),
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
for ri, row in enumerate(seg_data):
    bg = WHITE if ri%2==0 else RGBColor(0xE0,0xE8,0xF8)
    y = Inches(2.3 + ri*0.29)
    for ci, val in enumerate(row):
        add_rect(s, Inches(cx2[ci]), y, Inches(cw[ci]), Inches(0.28), fill=bg)
        fc = RGBColor(0x00,0x50,0xBB) if ci==1 else BLACK
        add_text(s, val, Inches(cx2[ci]+0.02), y+Inches(0.02),
                 Inches(cw[ci]-0.04), Inches(0.24), size=10.5, color=fc,
                 align=PP_ALIGN.CENTER, font="Courier New" if ci<=2 else "Calibri")

# Sơ đồ 7-seg
add_rect(s, Inches(0.3), Inches(4.85), Inches(6.3), Inches(2.42), fill=LGRAY)
add_textbox_multi(s, [
    ("Bit mapping:  bit0=SEG_A · bit1=SEG_B · · · bit6=SEG_G · bit7=DP", 13, False, NAVY),
    ("COM thứ tự:  COM0=chữ số giờ chục · COM1=giờ đơn · COM2=phút chục · COM3=phút đơn", 13, False, BLACK),
    ("Blink logic:  blink_flag 0/1 đổi mỗi 500ms → xóa segment_buff[digit] = 0x00", 13, False, BLUE),
    ("Ghost fix:   Tắt SEG trước → Tắt COM → Chuyển COM → Bật COM → Bật SEG", 13, True, GREEN),
], Inches(0.4), Inches(4.9), Inches(6.1), Inches(2.3))

add_rect(s, Inches(6.8), Inches(4.85), Inches(6.2), Inches(2.42), fill=LGRAY)
add_textbox_multi(s, [
    ("Tần số quét:  4 COM × 1ms = 4ms/chu kỳ → 250Hz", 13, True, BLUE),
    ("Mắt người:   Nhận biết flicker < 50Hz → 250Hz mượt hoàn toàn", 13, False, BLACK),
    ("update_display(): chọn data từ trạng thái hiện tại:", 13, False, NAVY),
    ("  NORMAL → g_hour/g_min", 12, False, BLACK),
    ("  SET_HOUR → g_hour nhấp nháy, g_min tĩnh", 12, False, BLACK),
    ("  SET_ALARM_* → alarm_hour/alarm_min, LED D6 nháy", 12, False, BLACK),
], Inches(6.9), Inches(4.9), Inches(6.0), Inches(2.3))

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — I2C & EEPROM
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "I2C & EEPROM AT24C02 — Lưu Báo Thức",
           "Dữ liệu không mất khi mất điện — giao thức I2C chuẩn")
slide_number(s, 9, TOTAL)

code_box(s, """// Ghi 1 byte vào EEPROM
void eeprom_write(uint8_t reg, uint8_t data) {
    I2C_Start();          // START condition
    I2C_WriteByte(0xA0);  // Slave address WRITE (7-bit: 0b1010_000 + R/W=0)
    I2C_WriteByte(reg);   // Địa chỉ thanh ghi (0x00=alarm_hour, 0x01=alarm_min)
    I2C_WriteByte(data);  // Dữ liệu cần lưu
    I2C_Stop();           // STOP condition
}

// Đọc 1 byte từ EEPROM (Repeated START)
uint8_t eeprom_read(uint8_t reg) {
    I2C_Start();          // START
    I2C_WriteByte(0xA0);  // Slave address WRITE
    I2C_WriteByte(reg);   // Chọn thanh ghi
    I2C_Start();          // Repeated START (không STOP giữa chừng!)
    I2C_WriteByte(0xA1);  // Slave address READ
    uint8_t val = I2C_ReadByte(NACK); // Đọc → NACK (byte cuối)
    I2C_Stop();
    return val;
}""",
         Inches(0.3), Inches(1.42), Inches(7.0), Inches(4.5), title="EEPROM Driver (I2C0)")

# Info bên phải
bullet_box(s, [
    ("AT24C02: 256 byte, endurance 1,000,000 lần ghi", NAVY, False),
    ("Địa chỉ slave: 0xA0 (write) / 0xA1 (read)", BLUE, False),
    ("EEPROM_AL_HOUR = 0x00 | EEPROM_AL_MIN = 0x01", BLUE, False),
    ("Kiểm tra 0xFF khi khởi động — phát hiện EEPROM mới chưa ghi", GREEN, True),
    ("Tuổi thọ: 1 write/ngày → ~1,370 NĂM", GREEN, True),
    ("PFPA: P0.10=SDA, P0.11=SCL — gán qua thanh ghi PFPA0", NAVY, False),
    ("Repeated START: tránh bus bị giải phóng giữa write và read", ORANGE, False),
],
           Inches(7.6), Inches(1.42), Inches(5.4), Inches(4.5),
           title="Thông số & Điểm quan trọng", title_color=NAVY, size=13)

# Timeline I2C
add_rect(s, Inches(0.3), Inches(6.05), Inches(12.7), Inches(1.2), fill=LGRAY)
add_textbox_multi(s, [
    ("Giao thức I2C Write:   START → 0xA0 → reg_addr → data → STOP", 13, False, NAVY),
    ("Giao thức I2C Read:    START → 0xA0 → reg_addr → rSTART → 0xA1 → data(NACK) → STOP", 13, False, BLUE),
], Inches(0.5), Inches(6.1), Inches(12.3), Inches(1.1))

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Buzzer PWM
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "Buzzer PWM — Âm Thanh Bằng Phần Cứng",
           "CT16B0 tạo sóng vuông 600Hz — không tốn CPU")
slide_number(s, 10, TOTAL)

code_box(s, """// CT16B0.c — Cấu hình PWM
SN_CT16B0->MR9     = 19999;      // Period: 12MHz/(19999+1) = 600Hz
SN_CT16B0->MR0     = 0;          // Duty=0 (tắt buzzer ban đầu)
SN_CT16B0->PWMCTRL = (1<<0)      // PWM0EN
                   | (1<<4)      // PWM0MODE (toggle)
                   | (1<<20);    // MR9RST (auto-reload)

// Bật buzzer
void buzzer_on(uint32_t freq_hz) {
    SN_CT16B0->MR9 = (12000000 / freq_hz) - 1; // Tính lại period
    SN_CT16B0->MR0 = SN_CT16B0->MR9 / 2;       // Duty 50%
}

// Tắt buzzer
void buzzer_off(void) {
    SN_CT16B0->MR0 = 0;          // Duty=0 → output thấp → im lặng
}""",
         Inches(0.3), Inches(1.42), Inches(7.0), Inches(4.35), title="CT16B0 PWM Driver")

bullet_box(s, [
    ("Pip 300ms: xác nhận nhấn phím — buzzer_on(600) → delay 300ms → buzzer_off()", NAVY, False),
    ("Alarm 5s: còi kêu 5 giây, 500ms on/off xen kẽ", BLUE, False),
    ("Priority: Alarm > Pip (alarm không bị pip override)", ORANGE, True),
    ("Output pin: P3.0 — gán qua PFPA3 → CT16B0 PWM0", NAVY, False),
    ("Ưu điểm PWM phần cứng: tần số chính xác tuyệt đối,\n  không bị ảnh hưởng bởi ngắt hay load CPU", GREEN, True),
    ("Verify: CT16B0→MR9=19999 → 12M÷20000=600Hz ✓", GREEN, True),
    ("Duty 50%: MR0=9999 = MR9÷2 → sóng vuông đối xứng", BLUE, False),
],
           Inches(7.6), Inches(1.42), Inches(5.4), Inches(4.35),
           title="Chế độ hoạt động", title_color=NAVY, size=13)

add_rect(s, Inches(0.3), Inches(5.9), Inches(12.7), Inches(1.35), fill=LGRAY)
add_textbox_multi(s, [
    ("Sóng PWM 600Hz, duty 50% (MR0=9999, MR9=19999, HCLK=12MHz):", 13, True, NAVY),
    ("│▔▔▔▔▔│_____│▔▔▔▔▔│_____│▔▔▔▔▔│_____│   ← 1 chu kỳ = 1/600Hz ≈ 1.67ms", 12, False, BLUE),
    ("MR0 match → toggle output LOW   |   MR9 match → toggle output HIGH + reset counter", 12, False, BLACK),
], Inches(0.45), Inches(5.95), Inches(12.3), Inches(1.25), wrap=False)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Watchdog & An toàn
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "Watchdog Timer — Bảo Vệ Hệ Thống",
           "Tự động reset nếu firmware treo — cơ chế quan trọng trong embedded thực tế")
slide_number(s, 11, TOTAL)

code_box(s, """// main() — while(1) loop
while (1) {
    /* ① Feed WDT — PHẢI làm đầu tiên mỗi vòng lặp */
    SN_WDT->FEED = 0x5AFA5AFA;   // Magic word để reset bộ đếm WDT

    /* ② Chờ tick 1ms */
    if (timer_1ms_flag) {
        timer_1ms_flag = 0;

        /* ③ Xử lý logic */
        Digital_Scan();           // Quét LED
        Key_Scan();               // Đọc phím
        process_key();            // Cập nhật state
        update_display();         // Cập nhật buffer LED
        ms_counter++;
        if (ms_counter >= 1000) {
            ms_counter = 0;
            tick_second();        // Đếm giây
        }
        update_buzzer();          // Cập nhật còi
    }
}
// Nếu vòng lặp bị treo → WDT không được feed → timeout ~250ms → RESET MCU""",
         Inches(0.3), Inches(1.42), Inches(7.2), Inches(5.3), title="Vòng lặp chính với WDT")

bullet_box(s, [
    ("WDT timeout: ~250ms nếu không feed", NAVY, False),
    ("Feed magic word: 0x5AFA5AFA — giá trị cố định của SN32F407", BLUE, False),
    ("Trường hợp WDT kích hoạt reset:", ORANGE, True),
    ("  • I2C busy-wait không thoát (EEPROM không phản hồi)", BLACK, False),
    ("  • Vòng lặp vô tận do bug logic", BLACK, False),
    ("  • Stack overflow gây HardFault", BLACK, False),
    ("Sau reset: MCU khởi động lại từ đầu, tải báo thức từ EEPROM", GREEN, True),
    ("Timeout SET mode: 30 giây không tương tác → NORMAL", BLUE, False),
    ("Volatile flag: đảm bảo compiler không optimize biến dùng chung ISR↔main", ORANGE, True),
],
           Inches(7.7), Inches(1.42), Inches(5.3), Inches(5.3),
           title="Phân tích an toàn hệ thống", title_color=NAVY, size=12.5)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Demo Keil Simulator
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "Demo — Keil MDK Simulator",
           "Mô phỏng hoàn toàn phần cứng, không cần board mạch thật")
slide_number(s, 12, TOTAL)

add_rect(s, Inches(0.3), Inches(1.42), Inches(4.8), Inches(5.82), fill=LGRAY)
add_text(s, "📋 Quy trình Demo", Inches(0.4), Inches(1.5), Inches(4.6), Inches(0.45),
         size=16, bold=True, color=NAVY)

steps = [
    ("1", "Build project (F7) → 0 Error"),
    ("2", "Debug Options → Use Simulator\nDLL: DARMCM1.DLL  param: -pCM0"),
    ("3", "Ctrl+F5 vào Debug mode"),
    ("4", "Mở Watch Window 1:\n  g_hour, g_min, g_sec, g_sim_key"),
    ("5", "F5 → chương trình chạy"),
    ("6", "Quan sát g_sec tăng → g_min cuộn"),
    ("7", "Set g_sim_key=0x14 → SET_HOUR mode"),
    ("8", "Set g_sim_key=0x22/0x42 → tăng/giảm"),
    ("9", "Cài báo thức: g_sim_key=0x88"),
]
for i, (num, step) in enumerate(steps):
    y = Inches(2.0 + i*0.58)
    add_rect(s, Inches(0.35), y, Inches(0.38), Inches(0.48), fill=BLUE)
    add_text(s, num, Inches(0.35), y+Inches(0.04), Inches(0.38), Inches(0.4),
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, step, Inches(0.8), y+Inches(0.04), Inches(4.25), Inches(0.5),
             size=12, color=BLACK)

# Bảng phím
add_rect(s, Inches(5.4), Inches(1.42), Inches(7.65), Inches(2.5), fill=LGRAY)
add_text(s, "🎮 Giả lập phím (Watch Window → g_sim_key)", Inches(5.5), Inches(1.49),
         Inches(7.4), Inches(0.42), size=15, bold=True, color=NAVY)

key_data = [
    ("SW3  — SETUP",  "0x14", "Vào/thoát SET mode"),
    ("SW6  — PLUS",   "0x22", "Tăng giá trị"),
    ("SW10 — MINUS",  "0x42", "Giảm giá trị"),
    ("SW16 — ALARM",  "0x88", "Vào SET_ALARM mode"),
]
kw = [2.1, 0.9, 3.5]
kx = [5.45, 7.56, 8.48]
for ci, h in enumerate(["Phím", "g_sim_key", "Chức năng"]):
    add_rect(s, Inches(kx[ci]), Inches(1.95), Inches(kw[ci]), Inches(0.36), fill=NAVY)
    add_text(s, h, Inches(kx[ci]), Inches(1.97), Inches(kw[ci]), Inches(0.32),
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
for ri, (k, v, desc) in enumerate(key_data):
    bg = WHITE if ri%2==0 else RGBColor(0xE0,0xE8,0xF8)
    y = Inches(2.33 + ri*0.37)
    for ci, val in enumerate([k, v, desc]):
        add_rect(s, Inches(kx[ci]), y, Inches(kw[ci]), Inches(0.35), fill=bg)
        fc = BLUE if ci==1 else BLACK
        fn = "Courier New" if ci==1 else "Calibri"
        add_text(s, val, Inches(kx[ci]+0.05), y+Inches(0.04),
                 Inches(kw[ci]-0.1), Inches(0.27), size=12, color=fc,
                 font=fn, align=PP_ALIGN.CENTER if ci==1 else PP_ALIGN.LEFT)

# Watch window mô phỏng
add_rect(s, Inches(5.4), Inches(4.05), Inches(7.65), Inches(3.2),
         fill=RGBColor(0x1E,0x1E,0x2E), line=CYAN, line_w=Pt(1).emu)
add_text(s, "Watch 1 — Keil Simulator", Inches(5.5), Inches(4.1), Inches(7.0), Inches(0.35),
         size=13, bold=True, color=CYAN)

watch_vars = [
    ("g_hour",         "uint8_t",  "09"),
    ("g_min",          "uint8_t",  "30"),
    ("g_sec",          "uint8_t",  "42"),
    ("g_sim_key",      "uint8_t",  "0x00"),
    ("alarm_hour",     "uint8_t",  "07"),
    ("alarm_min",      "uint8_t",  "00"),
    ("timer_1ms_flag", "volatile", "0"),
    ("state",          "uint8_t",  "0  (NORMAL)"),
]
add_textbox_multi(s, [
    (f"  {'Name':<18} {'Type':<12} Value", 11, True, CYAN),
    ("  " + "─"*45, 10, False, DGRAY),
] + [
    (f"  {n:<18} {t:<12} {v}", 11, False, RGBColor(0xD4,0xD4,0xD4))
    for n, t, v in watch_vars
], Inches(5.5), Inches(4.48), Inches(7.2), Inches(2.7), wrap=False)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Kết quả kiểm thử
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "Kết Quả Kiểm Thử", "14 test case — bao gồm các trường hợp biên")
slide_number(s, 13, TOTAL)

test_cases = [
    ("TC01", "Đếm giờ bình thường",         "g_sec: 0→59→0, g_min++",         "✅ PASS"),
    ("TC02", "Wrap-around 23:59→00:00",      "g_hour 23→0 khi g_min=59→0",     "✅ PASS"),
    ("TC03", "Cài giờ PLUS từ 0→23",         "g_hour tăng đúng, wrap 23→0",    "✅ PASS"),
    ("TC04", "Cài phút MINUS từ 0→59",       "g_min giảm, wrap 0→59",          "✅ PASS"),
    ("TC05", "Timeout 30s trong SET_HOUR",   "State → NORMAL sau 30000ms",     "✅ PASS"),
    ("TC06", "Cài báo thức và lưu EEPROM",  "alarm_hour/min ghi 0x00/0x01",   "✅ PASS"),
    ("TC07", "Báo thức kích hoạt",           "g_hour=alarm_hour, g_min=alarm_min, g_sec=0", "✅ PASS"),
    ("TC08", "Tắt báo thức bằng phím",       "buzzer_state → IDLE",            "✅ PASS"),
    ("TC09", "Pip 300ms khi nhấn phím",      "buzzer kêu 300ms rồi tắt",       "✅ PASS"),
    ("TC10", "LED D0 nhấp nháy 1Hz",         "P3.8 toggle mỗi 500ms",          "✅ PASS"),
    ("TC11", "Blink chữ số đang SET",        "blink_flag: 0/1 mỗi 500ms",      "✅ PASS"),
    ("TC12", "EEPROM 0xFF → dùng default",   "alarm=00:00 khi EEPROM mới",     "✅ PASS"),
    ("TC13", "WDT feed đều đặn",             "SN_WDT->FEED mỗi vòng lặp",     "✅ PASS"),
    ("TC14", "Không race condition",          "ISR chỉ set uint8_t flag",       "✅ PASS"),
]

hdr = ["ID", "Test Case", "Kết quả quan sát", "Verdict"]
hw = [0.7, 3.6, 4.5, 1.15]
hx = [0.3, 1.02, 4.64, 9.21]

for ci, h in enumerate(hdr):
    add_rect(s, Inches(hx[ci]), Inches(1.42), Inches(hw[ci]), Inches(0.4), fill=NAVY)
    add_text(s, h, Inches(hx[ci]), Inches(1.44), Inches(hw[ci]), Inches(0.36),
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for ri, (tid, name, obs, verdict) in enumerate(test_cases):
    bg = WHITE if ri%2==0 else LGRAY
    y = Inches(1.84 + ri*0.358)
    for ci, val in enumerate([tid, name, obs, verdict]):
        add_rect(s, Inches(hx[ci]), y, Inches(hw[ci]), Inches(0.35), fill=bg)
        fc = GREEN if "PASS" in val else (BLUE if ci==0 else BLACK)
        bold = "PASS" in val or ci==0
        add_text(s, val, Inches(hx[ci]+0.05), y+Inches(0.03),
                 Inches(hw[ci]-0.1), Inches(0.29),
                 size=11, color=fc, bold=bold,
                 align=PP_ALIGN.CENTER if ci in [0,3] else PP_ALIGN.LEFT)

add_text(s, "14/14 TEST CASES PASSED  ✅",
         Inches(0.3), Inches(7.0), Inches(13), Inches(0.4),
         size=16, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Đánh giá
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "Đánh Giá Chương Trình", "Điểm mạnh, điểm cần cải thiện, kết luận")
slide_number(s, 14, TOTAL)

bullet_box(s, [
    ("Register-level programming — không dùng HAL hay Arduino abstraction", GREEN, True),
    ("ISR tối giản: chỉ set flag, logic xử lý ở main loop", GREEN, True),
    ("volatile đúng chỗ — tránh compiler optimize mất cờ ISR", GREEN, False),
    ("GPIO atomic BSET/BCLR — thread-safe, không race condition", GREEN, False),
    ("PWM buzzer phần cứng — 600Hz ổn định, không tốn CPU", GREEN, False),
    ("EEPROM: kiểm tra 0xFF, tuổi thọ ~1,370 năm ở 1 write/ngày", GREEN, False),
    ("Ghost segment fix: tắt SEG trước khi chuyển COM", GREEN, False),
    ("WDT + timeout bảo vệ — sản phẩm không bao giờ bị treo", GREEN, False),
],
           Inches(0.3), Inches(1.42), Inches(6.15), Inches(4.2),
           title="✅ Điểm mạnh", title_color=GREEN, bullet="✓", size=13.5)

bullet_box(s, [
    ("IHRC ±5% drift: ~3 phút/giờ sai số (cần RTC ngoài cho sản phẩm thật)", ORANGE, True),
    ("I2C busy-wait: không có timeout → WDT sẽ reset nếu EEPROM không phản hồi", ORANGE, False),
    ("Chưa có __WFI() tiết kiệm điện giữa các ngắt", ORANGE, False),
    ("Debounce chỉ dựa polling period 1ms (chưa có bộ đếm 20ms)", ORANGE, False),
],
           Inches(0.3), Inches(5.72), Inches(6.15), Inches(1.52),
           title="⚠️  Cần cải thiện (nếu nâng cấp)", title_color=ORANGE, bullet="◦", size=13.5)

# Scorecard
add_rect(s, Inches(6.6), Inches(1.42), Inches(6.5), Inches(5.82), fill=LGRAY)
add_text(s, "📊 Bảng điểm", Inches(6.7), Inches(1.5), Inches(6.2), Inches(0.42),
         size=16, bold=True, color=NAVY)

scores = [
    ("Kiến trúc phần mềm",    "9/10", GREEN),
    ("Sử dụng phần cứng",     "9/10", GREEN),
    ("Reliability & Safety",  "7/10", ORANGE),
    ("Tính năng đề bài",      "10/10", GREEN),
    ("Khả năng mở rộng",      "8/10", BLUE),
    ("Tài liệu & Demo",       "10/10", GREEN),
]
for i, (crit, score, color) in enumerate(scores):
    y = Inches(2.0 + i*0.73)
    add_rect(s, Inches(6.65), y, Inches(4.5), Inches(0.62), fill=WHITE)
    add_text(s, crit, Inches(6.75), y+Inches(0.1), Inches(3.3), Inches(0.42),
             size=13.5, color=BLACK)
    add_rect(s, Inches(11.3), y, Inches(1.6), Inches(0.62), fill=color)
    add_text(s, score, Inches(11.3), y+Inches(0.1), Inches(1.6), Inches(0.42),
             size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(s, Inches(6.65), Inches(6.42), Inches(6.28), Inches(0.75), fill=NAVY)
add_text(s, "TỔNG:   8.8 / 10  ★★★★☆   Grade A",
         Inches(6.7), Inches(6.48), Inches(6.2), Inches(0.62),
         size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Hỏi & Đáp mẫu
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "Câu Hỏi Phỏng Vấn Thường Gặp", "Top 10 câu hỏi về dự án này")
slide_number(s, 15, TOTAL)

qas = [
    ("Tại sao dùng BSET/BCLR thay vì đọc-ghi-trả?",
     "BSET/BCLR là atomic — không thể bị ngắt giữa chừng. Read-Modify-Write gồm 3 lệnh riêng, ISR có thể chen vào giữa gây race condition."),
    ("SysTick RELOAD tính thế nào?",
     "RELOAD = HCLK/Freq − 1 = 12,000,000/1,000 − 1 = 11,999. Bộ đếm giảm từ 11999 về 0 mất đúng 1ms."),
    ("Tại sao ISR chỉ set flag, không xử lý trực tiếp?",
     "ISR cần ngắn gọn để không chặn các ngắt khác. Logic nặng (scan LED, I2C) ở main loop an toàn hơn."),
    ("volatile có tác dụng gì?",
     "Báo compiler: 'biến này có thể thay đổi bên ngoài luồng code' → không cache vào register → luôn đọc từ RAM."),
    ("Tại sao tắt SEG trước khi chuyển COM?",
     "Nếu chuyển COM khi SEG đang bật → segment sẽ sáng sai 1 khoảnh khắc ở digit mới (ghost). Tắt SEG trước loại bỏ hoàn toàn hiện tượng này."),
    ("I2C Repeated START khác STOP+START thế nào?",
     "Repeated START giữ quyền làm chủ bus (master). STOP+START giải phóng bus → slave khác có thể chiếm → nguy cơ đọc sai dữ liệu."),
    ("WDT timeout ~250ms được tính thế nào?",
     "Dựa vào giá trị thanh ghi WDT prescaler và clock WDT. Với SN32F407, default WDT clock từ ILRC ~32kHz, prescaler cho timeout ~250ms."),
    ("IHRC sai bao nhiêu? Ảnh hưởng đến đồng hồ thế nào?",
     "IHRC sai ±5% theo nhiệt độ. 5% của 3600s/giờ = 180s = 3 phút/giờ. Với sản phẩm thật cần thạch anh 32.768kHz ngoài + RTC."),
    ("State machine có bao nhiêu state? Kể tên.",
     "5 state: NORMAL(0), SET_HOUR(1), SET_MIN(2), SET_ALARM_HOUR(3), SET_ALARM_MIN(4)."),
    ("Buzzer 600Hz được tạo thế nào?",
     "CT16B0 MR9=19999: f = 12MHz/(19999+1) = 600Hz. MR0=9999 = MR9/2 → duty 50%. Output ra P3.0 qua PFPA."),
]

for i, (q, a) in enumerate(qas):
    col = 0 if i < 5 else 1
    row = i if i < 5 else i-5
    x = Inches(0.3) if col==0 else Inches(6.75)
    y = Inches(1.42 + row*1.1)
    w = Inches(6.25)
    add_rect(s, x, y, w, Inches(1.05), fill=LGRAY)
    add_text(s, f"Q{i+1}: {q}", x+Inches(0.1), y+Inches(0.04),
             w-Inches(0.2), Inches(0.38), size=12, bold=True, color=NAVY)
    add_text(s, f"→ {a}", x+Inches(0.1), y+Inches(0.42),
             w-Inches(0.2), Inches(0.6), size=11, color=BLACK)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Cấu trúc file
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "Cấu Trúc Project", "Tổ chức source code theo chức năng")
slide_number(s, 16, TOTAL)

code_box(s, """Clock_SN32F407/
├── Source/
│   ├── UserAPP/
│   │   ├── main.c          ← Vòng lặp chính, state machine, logic nghiệp vụ
│   │   └── main_sim.c      ← Phiên bản Simulator (SysTick thay CT16B1)
│   ├── Driver/
│   │   ├── CT16B0.c/.h     ← Timer buzzer PWM 600Hz
│   │   ├── CT16B1.c/.h     ← Timer 1ms (phần cứng)
│   │   ├── GPIO.c/.h       ← Cấu hình tất cả GPIO pin
│   │   └── SysTick_1ms.c/.h← Timer 1ms (simulator)
│   ├── Module/
│   │   ├── Segment.c/.h    ← SEGMENT_TABLE, Digital_Scan()
│   │   └── EEPROM.c/.h     ← I2C driver + AT24C02 read/write
│   └── CMSIS/
│       └── SN32F4xx/       ← Device header, startup_SN32F407.s
└── Clock_SN32F407.uvprojx  ← Keil µVision5 project file""",
         Inches(0.3), Inches(1.42), Inches(6.5), Inches(5.5), title="Cây thư mục")

bullet_box(s, [
    ("main.c: ~900 dòng — toàn bộ logic ứng dụng", NAVY, True),
    ("CT16B1.c: timer 1ms cho phần cứng thật", BLUE, False),
    ("SysTick_1ms.c: timer 1ms cho Keil Simulator", BLUE, False),
    ("CT16B0.c: PWM buzzer — tần số lập trình được", BLUE, False),
    ("GPIO.c: cấu hình P0(SEG), P1(COM), P3(LED+Buzzer)", BLUE, False),
    ("Segment.c: bảng mã 7 đoạn + hàm quét multiplexing", BLUE, False),
    ("EEPROM.c: giao tiếp I2C với AT24C02", BLUE, False),
    ("startup_SN32F407.s: assembly — reset handler, vector table", ORANGE, False),
    ("Build output: Clock_SN32F407.hex → flash lên board thật", GREEN, True),
],
           Inches(7.0), Inches(1.42), Inches(6.1), Inches(5.5),
           title="Vai trò từng file", title_color=NAVY, size=13.5)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Tài liệu đi kèm
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "Tài Liệu Đi Kèm", "Bộ tài liệu đầy đủ cho dự án")
slide_number(s, 17, TOTAL)

docs = [
    ("📄", "Báo_cáo_Digital_Clock_SN32F407.docx",
     "Báo cáo chính thức — mô tả yêu cầu, thiết kế, kết quả\nKèm Phụ lục: Edge cases, Race condition, EEPROM wear", BLUE),
    ("📚", "Tài_liệu_giải_thích_code.docx",
     "Giải thích chi tiết từng dòng code — 16 phần\nDành cho đọc hiểu sâu trước báo cáo / phỏng vấn", GREEN),
    ("🎓", "Tài_liệu_phỏng_vấn_SN32F407.docx",
     "13 phần lý thuyết + 20 Q&A phỏng vấn\nCover: MCU, GPIO, Timer, I2C, WDT, State machine, IHRC drift...", ORANGE),
    ("🎬", "Kịch_bản_quay_demo.docx",
     "12 cảnh quay trên Keil Simulator\nKèm: Lời dẫn voiceover, hướng dẫn Keil setup, bảng phím", RED),
    ("🖥", "Slide_Digital_Clock_SN32F407.pptx",
     "Slide thuyết trình 18 trang — tài liệu này\nDùng cho bảo vệ đồ án / trình bày nhóm", NAVY),
]

for i, (icon, name, desc, color) in enumerate(docs):
    y = Inches(1.5 + i*1.05)
    add_rect(s, Inches(0.3), y, Inches(12.7), Inches(0.98), fill=LGRAY)
    add_rect(s, Inches(0.3), y, Inches(0.55), Inches(0.98), fill=color)
    add_text(s, icon, Inches(0.3), y+Inches(0.2), Inches(0.55), Inches(0.58),
             size=22, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, name, Inches(0.95), y+Inches(0.06), Inches(11.9), Inches(0.4),
             size=14, bold=True, color=color)
    add_text(s, desc, Inches(0.98), y+Inches(0.48), Inches(11.7), Inches(0.46),
             size=12, color=DGRAY)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Kết luận
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=NAVY)
add_rect(s, 0, 0, W, Inches(0.18), fill=CYAN)
add_rect(s, 0, H-Inches(0.18), W, Inches(0.18), fill=CYAN)
add_rect(s, 0, Inches(2.5), W, Pt(3).emu, fill=CYAN)

add_text(s, "Kết Luận", Inches(0.4), Inches(0.5), Inches(12.5), Inches(0.85),
         size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Digital Clock — SN32F407 ARM Cortex-M0",
         Inches(0.4), Inches(1.35), Inches(12.5), Inches(0.55),
         size=20, italic=True, color=CYAN, align=PP_ALIGN.CENTER)

summary_lines = [
    "✅  15/15 yêu cầu đề bài đã hoàn thành",
    "✅  Register-level programming — hiểu sâu phần cứng ARM Cortex-M0",
    "✅  State machine 5 trạng thái — kiến trúc rõ ràng, dễ mở rộng",
    "✅  ISR chuẩn + volatile + atomic GPIO — code đúng chuẩn embedded",
    "✅  WDT + timeout — hệ thống an toàn, không bao giờ bị treo",
    "✅  EEPROM persistence — báo thức không mất khi mất điện",
    "✅  Demo hoàn toàn trên Keil Simulator — không cần board thật",
]
tb = s.shapes.add_textbox(Inches(1.2), Inches(2.75), Inches(10.9), Inches(4.0))
tf = tb.text_frame; tf.word_wrap = True
for i, line in enumerate(summary_lines):
    p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
    p.space_before = Pt(4)
    run = p.add_run(); run.text = line
    run.font.name = "Calibri"; run.font.size = Pt(17)
    run.font.color.rgb = WHITE if "✅" not in line else RGBColor(0x7F,0xFF,0xAA)
    run.font.bold = True

add_text(s, "Cảm ơn thầy/cô và các bạn đã lắng nghe!",
         Inches(0.4), Inches(6.9), Inches(12.5), Inches(0.48),
         size=18, italic=True, color=CYAN, align=PP_ALIGN.CENTER)

slide_number(s, 18, TOTAL)

# ════════════════════════════════════════════════════════════════════════════
prs.save(OUT)
print("✅ Slide đã tạo:", OUT)
print(f"   {TOTAL} slides, 16:9, thiết kế Navy/White/Cyan")
